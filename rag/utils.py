# rag/utils.py
"""
Utilities to load documents, extract text + images, and return chunk tuples for indexing.

Each returned item is a tuple: (content: str, metadata: dict)

- Text chunk:
    content = "...text..."
    metadata = {
        "type": "text",
        "orig_filename": "<name>",
        "page": <int or None>
    }

- Image chunk:
    content = "[Figure <n>] <short caption>"
    metadata = {
        "type": "image",
        "image_path": "<relative path under uploads/>",
        "orig_filename": "<name>",
        "page": <int or None>,
        "figure_no": <int>,
        "caption": "<fuller caption if available>"
    }

Supported types: PDF, DOCX, PPTX, TXT/CSV/MD (text only)
Images are saved to: uploads/<docname>_assets/<file> and referenced via metadata.image_path
"""

from __future__ import annotations

import os
from typing import List, Tuple, Dict, Any

# ---- Optional deps (install latest) ----
# pip install --upgrade pymupdf python-docx python-pptx
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ------------ Chunking config ------------
CHUNK_SIZE = 900
CHUNK_OVERLAP = 150


# ------------- Helpers ------------------
def _ensure_dir(path: str):
    os.makedirs(path, exist_ok=True)


def _slug(s: str) -> str:
    return "".join(c if c.isalnum() or c in ("-", "_", ".") else "_" for c in s)


def _norm_ws(s: str) -> str:
    return " ".join((s or "").split())


def _chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Greedy fixed-size chunking with overlap."""
    text = _norm_ws(text)
    if not text:
        return []
    chunks: List[str] = []
    i = 0
    step = max(size - overlap, 1)
    while i < len(text):
        chunks.append(text[i:i + size])
        i += step
    return chunks


def _text_meta_chunk(text: str, orig_filename: str, page: int | None = None) -> Tuple[str, Dict[str, Any]]:
    meta = {"type": "text", "orig_filename": orig_filename, "page": page}
    return text, meta


def _image_meta_chunk(
    caption_short: str,
    rel_image_path: str,
    orig_filename: str,
    page: int | None,
    figure_no: int,
    caption_full: str | None = None,
) -> Tuple[str, Dict[str, Any]]:
    content = f"[Figure {figure_no}] {caption_short or 'Image from document'}"
    meta = {
        "type": "image",
        "image_path": rel_image_path,
        "orig_filename": orig_filename,
        "page": page,
        "figure_no": figure_no,
        "caption": caption_full or caption_short or "",
    }
    return content, meta


def _assets_dir(filepath: str, uploads_root: str) -> str:
    base = os.path.splitext(os.path.basename(filepath))[0]
    dir_path = os.path.join(uploads_root, f"{_slug(base)}_assets")
    _ensure_dir(dir_path)
    return dir_path


# ------------- PDF loader ----------------
def _load_pdf_with_images(filepath: str, uploads_root: str) -> List[Tuple[str, Dict[str, Any]]]:
    """
    Extract:
      - page text (chunked, with page numbers)
      - images (saved as PNGs under assets dir), with simple captions from nearby page text
    """
    results: List[Tuple[str, Dict[str, Any]]] = []
    doc = fitz.open(filepath)
    assets_dir = _assets_dir(filepath, uploads_root)
    filename = os.path.basename(filepath)

    figure_no = 1
    for pidx in range(len(doc)):
        page = doc.load_page(pidx)

        # Page text
        page_text = page.get_text("text") or ""
        for chunk in _chunk_text(page_text):
            results.append(_text_meta_chunk(chunk, filename, page=pidx + 1))

        # Images
        image_list = page.get_images(full=True)
        if not image_list:
            continue

        # Simple caption context from first lines of page
        lines = [ln.strip() for ln in (page_text or "").splitlines() if ln.strip()]
        simple_caption = _norm_ws(" ".join(lines[:4]))[:300]  # brief

        for idx, img in enumerate(image_list, start=1):
            try:
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                if pix.n > 4:  # CMYK → RGB
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                img_name = f"{_slug(os.path.splitext(filename)[0])}_p{pidx+1}_{idx}.png"
                save_abs = os.path.join(assets_dir, img_name)
                pix.save(save_abs)
                rel_path = os.path.relpath(save_abs, uploads_root).replace("\\", "/")
                results.append(
                    _image_meta_chunk(
                        caption_short=simple_caption,
                        rel_image_path=rel_path,
                        orig_filename=filename,
                        page=pidx + 1,
                        figure_no=figure_no,
                        caption_full=simple_caption,
                    )
                )
                figure_no += 1
            except Exception:
                # Non-fatal: skip broken image
                continue

    doc.close()
    return results


# ------------- DOCX loader ---------------
def _load_docx_with_images(filepath: str, uploads_root: str) -> List[Tuple[str, Dict[str, Any]]]:
    """
    Extract:
      - full text (chunked)
      - embedded images via package parts; captions approximated from nearby paragraphs
    """
    results: List[Tuple[str, Dict[str, Any]]] = []
    doc = Document(filepath)
    assets_dir = _assets_dir(filepath, uploads_root)
    filename = os.path.basename(filepath)

    # Text
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    big_text = "\n".join(paragraphs)
    for chunk in _chunk_text(big_text):
        results.append(_text_meta_chunk(chunk, filename, page=None))

    # Approx caption source
    context_snippet = _norm_ws(" ".join(paragraphs[:6]))[:300]

    # Embedded images (best-effort across python-docx versions)
    figure_no = 1
    image_parts = []
    try:
        # newer approach: iterate related parts that expose .blob
        rels = getattr(doc.part.package, "part_related_by", None)
        if rels:
            image_parts = [p for p in rels.values() if hasattr(p, "blob")]
    except Exception:
        pass

    # fallback: some versions expose image_parts attribute
    if not image_parts and hasattr(doc.part.package, "image_parts"):
        try:
            image_parts = list(doc.part.package.image_parts)
        except Exception:
            image_parts = []

    # Save images
    for i, part in enumerate(image_parts, start=1):
        try:
            ct = getattr(part, "content_type", "") or ""
            ext = ".png"
            if "jpeg" in ct or "jpg" in ct:
                ext = ".jpg"
            elif "png" in ct:
                ext = ".png"
            elif "gif" in ct:
                ext = ".gif"
            img_name = f"{_slug(os.path.splitext(filename)[0])}_fig_{i}{ext}"
            save_abs = os.path.join(assets_dir, img_name)
            with open(save_abs, "wb") as f:
                f.write(part.blob)
            rel_path = os.path.relpath(save_abs, uploads_root).replace("\\", "/")

            results.append(
                _image_meta_chunk(
                    caption_short=context_snippet,
                    rel_image_path=rel_path,
                    orig_filename=filename,
                    page=None,
                    figure_no=figure_no,
                    caption_full=context_snippet,
                )
            )
            figure_no += 1
        except Exception:
            continue

    return results


# ------------- PPTX loader ---------------
def _load_pptx_with_images(filepath: str, uploads_root: str) -> List[Tuple[str, Dict[str, Any]]]:
    """
    Extract:
      - slide text (chunked per deck)
      - picture shapes (images), with slide text used as caption
    """
    results: List[Tuple[str, Dict[str, Any]]] = []
    prs = Presentation(filepath)
    assets_dir = _assets_dir(filepath, uploads_root)
    filename = os.path.basename(filepath)

    figure_no = 1
    for sidx, slide in enumerate(prs.slides, start=1):
        # Slide text
        slide_texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                txt = (shape.text or "").strip()
                if txt:
                    slide_texts.append(txt)
            elif hasattr(shape, "text"):
                txt = (shape.text or "").strip()
                if txt:
                    slide_texts.append(txt)
        slide_text = "\n".join(slide_texts)
        for chunk in _chunk_text(slide_text):
            results.append(_text_meta_chunk(chunk, filename, page=sidx))

        # Images
        caption_ctx = _norm_ws(slide_text)[:300]
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    ext = os.path.splitext(img.filename)[1] or ".png"
                    img_name = f"{_slug(os.path.splitext(filename)[0])}_s{sidx}_{figure_no}{ext}"
                    save_abs = os.path.join(assets_dir, img_name)
                    with open(save_abs, "wb") as f:
                        f.write(img.blob)
                    rel_path = os.path.relpath(save_abs, uploads_root).replace("\\", "/")

                    # Prefer alt text as caption if present
                    alt = ""
                    try:
                        alt = (getattr(shape, "alternative_text", "") or "").strip()
                    except Exception:
                        alt = ""
                    caption = _norm_ws(alt) or caption_ctx

                    results.append(
                        _image_meta_chunk(
                            caption_short=caption,
                            rel_image_path=rel_path,
                            orig_filename=filename,
                            page=sidx,
                            figure_no=figure_no,
                            caption_full=caption,
                        )
                    )
                    figure_no += 1
            except Exception:
                continue

    return results


# --------- Plain text / CSV / MD ---------
def _load_plain(filepath: str, uploads_root: str) -> List[Tuple[str, Dict[str, Any]]]:
    results: List[Tuple[str, Dict[str, Any]]] = []
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
    except Exception:
        text = ""
    for chunk in _chunk_text(text):
        results.append(_text_meta_chunk(chunk, os.path.basename(filepath)))
    return results


# --------------- Public API ---------------
def load_and_split_with_images(filepath: str, uploads_root: str = "uploads") -> List[Tuple[str, Dict[str, Any]]]:
    """
    Detect file type and return list of (content, metadata) chunks.
    Image files saved to uploads/<doc>_assets, referenced by relative metadata.image_path.
    """
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return _load_pdf_with_images(filepath, uploads_root)
    if ext == ".docx":
        return _load_docx_with_images(filepath, uploads_root)
    if ext == ".pptx":
        return _load_pptx_with_images(filepath, uploads_root)
    if ext in (".txt", ".csv", ".md"):
        return _load_plain(filepath, uploads_root)
    # Fallback: treat as plain text
    return _load_plain(filepath, uploads_root)
